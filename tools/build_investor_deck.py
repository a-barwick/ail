#!/usr/bin/env python3
"""Build the AIL investor deck as a Google Slides-importable PowerPoint file.

Requires python-pptx 1.0 or newer:
    python3 -m pip install --user 'python-pptx>=1.0,<2'
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "docs" / "investor-deck.pptx"

W = Inches(13.333)
H = Inches(7.5)

NAVY = "09111F"
PANEL = "111D30"
PANEL_2 = "17263C"
WHITE = "F5F7FA"
MUTED = "AAB8CC"
CYAN = "4DD9E7"
BLUE = "5E8BFF"
GREEN = "61D095"
AMBER = "F5B942"
RED = "F07178"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color: str, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 20,
    color: str = WHITE,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    line_spacing: float | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    if line_spacing is not None:
        paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def rect(slide, x, y, w, h, fill=PANEL, line=PANEL, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def line(slide, x1, y1, x2, y2, color=MUTED, width=1.5, arrow=False):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    set_line(shape, color, width)
    if arrow:
        shape.line.end_arrowhead = True
    return shape


def dot(slide, x, y, d=0.12, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    set_fill(shape, color)
    set_line(shape, color)
    return shape


def pill(slide, x, y, w, text, color=CYAN, text_color=NAVY):
    shape = rect(slide, x, y, w, 0.34, color, color, radius=True)
    shape.adjustments[0] = 0.45
    textbox(
        slide, x, y + 0.01, w, 0.28, text.upper(), size=9, color=text_color,
        bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )
    return shape


def base_slide(prs, number: int, section: str = "INVESTOR DECK"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(NAVY)
    textbox(slide, 0.55, 0.25, 2.6, 0.25, f"AIL  /  {section}", size=9, color=MUTED, bold=True)
    textbox(slide, 12.25, 7.08, 0.5, 0.2, f"{number:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    return slide


def title(slide, text: str, sub: str | None = None):
    size = 27 if len(text) > 46 else 30
    textbox(slide, 0.62, 0.62, 12.0, 0.88, text, size=size, bold=True)
    if sub:
        textbox(slide, 0.64, 1.53, 11.6, 0.38, sub, size=13, color=MUTED)


def add_notes(slide, text: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    frame.text = text.strip()


def card(slide, x, y, w, h, heading, body, *, accent=CYAN, heading_size=16, body_size=12):
    rect(slide, x, y, w, h, PANEL, PANEL_2)
    rect(slide, x, y, 0.07, h, accent, accent, radius=False)
    textbox(slide, x + 0.22, y + 0.2, w - 0.38, 0.34, heading, size=heading_size, bold=True)
    textbox(slide, x + 0.22, y + 0.7, w - 0.38, h - 0.82, body, size=body_size, color=MUTED)


def metric(slide, x, y, w, value, label, *, color=CYAN):
    rect(slide, x, y, w, 1.3, PANEL, PANEL_2)
    textbox(slide, x + 0.18, y + 0.12, w - 0.36, 0.56, value, size=28, color=color, bold=True)
    textbox(slide, x + 0.18, y + 0.75, w - 0.36, 0.35, label, size=11, color=MUTED)


def slide_01(prs):
    s = base_slide(prs, 1, "THESIS")
    pill(s, 0.64, 0.82, 1.55, "Investor brief", CYAN)
    textbox(s, 0.64, 1.52, 9.5, 1.5, "A programming language for\nsoftware built by agents", size=39, bold=True)
    textbox(s, 0.68, 3.45, 8.3, 0.65, "AIL is an executable language and semantic compiler designed to reduce the work between generated code and a trusted change.", size=18, color=MUTED)
    # Visual: model -> control layer -> publish
    for x, label, c in [(0.72, "MODEL +\nAGENT", BLUE), (4.42, "AIL\nCOMPILER", CYAN), (8.12, "VALIDATED\nCANDIDATE", GREEN)]:
        rect(s, x, 5.05, 2.65, 1.15, PANEL, c)
        textbox(s, x, 5.25, 2.65, 0.68, label, size=15, color=c, bold=True, align=PP_ALIGN.CENTER)
    line(s, 3.42, 5.63, 4.25, 5.63, BLUE, 2.5, True)
    line(s, 7.12, 5.63, 7.95, 5.63, CYAN, 2.5, True)
    textbox(s, 10.95, 5.13, 1.55, 0.95, "Compiler\nruns today", size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "AIL is not another coding agent. It is an executable language plus a compiler interface that makes consequences, authority, validation, and publication mechanically inspectable. The compiler runs today. No AIL-versus-baseline agent comparison has run.")


def slide_02(prs):
    s = base_slide(prs, 2, "MARKET SHIFT")
    title(s, "Code generation is becoming abundant", "The scarce resource moves from producing code to trusting change.")
    # Diverging bars
    textbox(s, 0.78, 2.15, 2.5, 0.3, "RELATIVE WORK", size=9, color=MUTED, bold=True)
    line(s, 1.02, 5.95, 11.95, 5.95, MUTED, 1)
    line(s, 1.02, 2.48, 1.02, 5.95, MUTED, 1)
    line(s, 1.15, 3.05, 11.5, 5.35, BLUE, 5)
    line(s, 1.15, 5.35, 11.5, 2.75, CYAN, 5)
    pill(s, 8.92, 2.3, 2.45, "Trust + verification", CYAN)
    pill(s, 8.92, 5.18, 2.45, "First-pass generation", BLUE, WHITE)
    textbox(s, 0.9, 6.17, 3.2, 0.28, "PAST: HUMAN-SCARCE OUTPUT", size=9, color=MUTED, bold=True)
    textbox(s, 9.2, 6.17, 2.8, 0.28, "AGENT-SCALE OUTPUT", size=9, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    card(s, 9.2, 3.48, 2.7, 1.2, "New bottleneck", "Context • consequence • repair • regression control", accent=AMBER, body_size=11)
    textbox(s, 1.18, 2.53, 5.5, 0.26, "CONCEPTUAL — NO MEASURED MARKET OR WORKLOAD DATA", size=8.5, color=AMBER, bold=True)
    add_notes(s, "Foundation models change the supply curve for plausible code. That does not make software change cheap: the costly work shifts to finding the right context, understanding consequences, repairing failures, and proving that a change is safe. AIL is built around that shifted cost structure. This chart is conceptual, not measured market data.")


def slide_03(prs):
    s = base_slide(prs, 3, "PROBLEM")
    title(s, "Today’s agents repeatedly reconstruct the program", "Strong tools expose substantial semantics; agents still reconstruct parts of the change surface.")
    labels = ["Search files", "Read context", "Infer effects", "Edit", "Compile / test", "Repair"]
    xs = [0.68, 2.65, 4.62, 6.59, 8.56, 10.53]
    for i, (x, label) in enumerate(zip(xs, labels)):
        color = RED if label == "Repair" else PANEL_2
        rect(s, x, 2.25, 1.48, 0.82, color, color)
        textbox(s, x + 0.08, 2.48, 1.32, 0.28, label, size=11, bold=True, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            line(s, x + 1.5, 2.66, x + 1.9, 2.66, MUTED, 1.7, True)
    line(s, 11.27, 3.15, 11.27, 3.75, RED, 2)
    line(s, 11.27, 3.75, 1.45, 3.75, RED, 2)
    line(s, 1.45, 3.75, 1.45, 3.17, RED, 2, True)
    textbox(s, 4.38, 3.9, 4.7, 0.38, "RECONSTRUCTION LOOP", size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    for x, h, b in [
        (0.78, "Hidden consequence", "A schema edit can touch producers, storage, projections, tests, and evidence."),
        (4.45, "Hidden authority", "Passing tests can still move database access into the wrong layer."),
        (8.12, "Unbounded review", "Humans reconstruct what changed and what the agent may have missed."),
    ]:
        card(s, x, 4.62, 3.25, 1.6, h, b, accent=RED, body_size=11)
    add_notes(s, "Modern compilers, language servers, refactors, search, and tests already expose substantial semantics. Even with those normal strong tools, the agent often still assembles a file-oriented view and reasons across parts of it probabilistically. That reconstruction repeats after diagnostics and again during review. AIL must demonstrate a further reduction in total work against that baseline, not against raw text editing.")


def slide_04(prs):
    s = base_slide(prs, 4, "PRODUCT THESIS")
    title(s, "AIL exposes language guarantees through the compiler")
    rect(s, 0.68, 2.0, 12.0, 4.35, PANEL, PANEL_2)
    # center compiler
    rect(s, 4.72, 2.72, 3.9, 2.9, PANEL_2, CYAN)
    textbox(s, 5.1, 3.05, 3.15, 0.5, "SEMANTIC COMPILER", size=21, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 5.18, 3.78, 3.0, 1.15, "canonical source\n+ typed semantic model\n+ immutable revisions", size=13, color=WHITE, align=PP_ALIGN.CENTER)
    items = [
        (1.0, 2.45, "INSPECTION", "revision-scoped semantic facts"),
        (1.0, 4.55, "BOUNDED IMPACT", "must change / review / unchecked"),
        (9.2, 2.45, "AUTHORITY", "effects + capabilities"),
        (9.2, 4.55, "PUBLICATION", "validate + commit or roll back"),
    ]
    for x, y, h, b in items:
        rect(s, x, y, 2.65, 1.05, NAVY, PANEL_2)
        textbox(s, x + 0.18, y + 0.15, 2.3, 0.26, h, size=12, color=CYAN, bold=True)
        textbox(s, x + 0.18, y + 0.52, 2.3, 0.24, b, size=10, color=MUTED)
        if x < 4:
            line(s, x + 2.68, y + 0.52, 4.55, 4.18, CYAN, 1.5, True)
        else:
            line(s, 8.78, 4.18, x - 0.12, y + 0.52, CYAN, 1.5, True)
    add_notes(s, "The key design choice is that the language creates semantics the compiler can expose completely. Canonical text remains the durable, human-auditable source. The semantic model is the operational interface for agents: exact impact, explicit authority, deterministic validation, and atomic publication. This is more than an IDE wrapper because the guarantees originate in the language contract.")


def slide_05(prs):
    s = base_slide(prs, 5, "WORKFLOW")
    title(s, "From reconstruct-and-repair to query-change-prove")
    pill(s, 0.74, 1.82, 1.2, "Before", RED, WHITE)
    pill(s, 6.75, 1.82, 1.2, "With AIL", GREEN, NAVY)
    before = ["Search repository", "Infer dependencies", "Coordinate text edits", "Run broad checks", "Repair surprises"]
    after = ["Query supported impact", "Inspect authority + policy", "Submit whole candidate", "Compiler validates evidence", "Publish or roll back"]
    for col, items, accent in [(0.74, before, RED), (6.75, after, GREEN)]:
        for i, item in enumerate(items):
            y = 2.48 + i * 0.76
            dot(s, col + 0.05, y + 0.2, 0.13, accent)
            textbox(s, col + 0.36, y, 4.75, 0.5, item, size=16, bold=i in (0, 4))
            if i < 4:
                line(s, col + 0.115, y + 0.36, col + 0.115, y + 0.88, accent, 1.5)
    line(s, 6.23, 1.92, 6.23, 6.33, PANEL_2, 2)
    textbox(s, 4.99, 6.47, 2.5, 0.32, "LESS RECONSTRUCTION", size=10, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "The before-and-after is the investor demo’s central visual. AIL does not claim the agent stops using judgment. It moves complete, deterministic work out of model inference: identifying affected semantic roles, checking authority and architecture policy, validating the whole candidate, and publishing atomically. The economic hypothesis is that this shortens the loop and makes more changes safely autonomous.")


def slide_06(prs):
    s = base_slide(prs, 6, "PROOF TODAY")
    title(s, "The Rust compiler already validates complete changes", "Checked-in contracts and deterministic fixtures define the behavior.")
    stages = [
        ("SOURCE", "lossless parse\ncanonical format", BLUE),
        ("SEMANTICS", "types • effects\ncapabilities", CYAN),
        ("REVISION", "handles • inspect\nvalidated rename", CYAN),
        ("CHANGE", "impact • semantic diff\natomic transaction", GREEN),
        ("POLICY", "architecture delta\ncommit / rollback", GREEN),
    ]
    for i, (h, b, c) in enumerate(stages):
        x = 0.58 + i * 2.52
        rect(s, x, 2.15, 2.1, 1.55, PANEL, c)
        textbox(s, x + 0.15, 2.37, 1.8, 0.28, h, size=12, color=c, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, x + 0.15, 2.85, 1.8, 0.52, b, size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 4:
            line(s, x + 2.13, 2.93, x + 2.43, 2.93, c, 2, True)
    metric(s, 0.74, 4.48, 2.72, "37", "public behavior fixtures", color=GREEN)
    metric(s, 3.72, 4.48, 2.72, "12", "semantic relationship kinds", color=CYAN)
    metric(s, 6.70, 4.48, 2.72, "23", "architecture contract scenarios", color=CYAN)
    metric(s, 9.68, 4.48, 2.72, "4", "behavior-equivalent references", color=BLUE)
    textbox(s, 0.76, 6.2, 11.6, 0.45, "WORKS NOW: deterministic compiler enforcement  •  NOT RUN: AIL-versus-baseline agent comparison", size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "This is not a mock interface. The Rust compiler parses and checks canonical AIL, executes the reference service, stores immutable revisions, computes exact impact, validates whole-workspace schema evolution, enforces architecture policy atomically, and supports modules, import aliases, qualified references, and ordinary calls. The counts here are repository facts, not performance claims: 37 public behavior fixtures, 12 relationship kinds in the impact graph, 23 architecture scenarios, and four equivalent baseline implementations.")


def slide_07(prs):
    s = base_slide(prs, 7, "IMPACT PROOF")
    title(s, "Before editing, the compiler returns the bounded change surface", "Supported fixture: add required priority while preserving V1 compatibility.")
    metric(s, 0.72, 2.04, 2.7, "12", "must change", color=GREEN)
    metric(s, 3.62, 2.04, 2.7, "2", "review with reason", color=AMBER)
    metric(s, 6.52, 2.04, 2.7, "1", "unchecked boundary", color=RED)
    rect(s, 9.42, 2.04, 2.95, 1.3, PANEL, PANEL_2)
    textbox(s, 9.65, 2.22, 2.52, 0.33, "AUTHORITY", size=11, color=CYAN, bold=True)
    textbox(s, 9.65, 2.67, 2.52, 0.32, "unchanged", size=20, color=GREEN, bold=True)
    roles = [
        "request + stored schemas", "V1 adapters", "handler construction",
        "store capability", "persisted encoder", "V1 / V2 projections",
        "fixture + completion evidence",
    ]
    rect(s, 0.72, 3.74, 11.65, 2.32, PANEL, PANEL_2)
    textbox(s, 0.98, 3.96, 3.1, 0.34, "COMPILER-DERIVED ROLES", size=11, color=CYAN, bold=True)
    for i, role in enumerate(roles):
        col = i % 4
        row = i // 4
        x = 1.0 + col * 2.85
        y = 4.55 + row * 0.67
        dot(s, x, y + 0.08, 0.12, GREEN)
        textbox(s, x + 0.25, y, 2.45, 0.34, role, size=11, color=WHITE)
    textbox(s, 0.76, 6.31, 11.6, 0.34, "Honest coverage: external clients stay visible as unchecked—not silently omitted.", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(s, "The demo should show this as a visual work list, not raw JSON. For the locked priority-evolution fixture, the compiler derives exactly 12 must-change locations, two review sites with reasons, and one known external boundary whose source is unavailable. The successful transaction uses five ordered whole-path edits. It also proves capabilities, effects, and ordering are unchanged. These counts are fixture-specific proof, not a claim about all programs.")


def slide_08(prs):
    s = base_slide(prs, 8, "POLICY PROOF")
    title(s, "Passing behavior evidence is necessary. It is not sufficient.", "Three frozen semantic candidate fixtures carry the same accepted 6/6 behavior evidence.")
    candidates = [
        (0.7, "DOMAIN-OWNED", "6 / 6 behavior", "0 findings", "PUBLISHED", GREEN),
        (4.47, "CENTRALIZED", "6 / 6 behavior", "4 findings", "ROLLED BACK", RED),
        (8.24, "HELPER-SPLIT", "6 / 6 behavior", "3 findings", "ROLLED BACK", RED),
    ]
    for x, name, behavior, findings, result, c in candidates:
        rect(s, x, 2.18, 3.42, 3.55, PANEL, c)
        pill(s, x + 0.24, 2.43, 1.62, name, c, NAVY if c == GREEN else WHITE)
        textbox(s, x + 0.26, 3.18, 2.9, 0.38, behavior, size=18, bold=True)
        textbox(s, x + 0.26, 3.75, 2.9, 0.34, findings, size=14, color=c, bold=True)
        if name == "CENTRALIZED":
            body = "dispatch growth\nstore authority in transport\nstate + dependency boundary"
        elif name == "HELPER-SPLIT":
            body = "small helpers individually\naggregate transport boundary\nstill violated"
        else:
            body = "authority remains in domain\nexisting dispatch debt unchanged\ncomplete evidence"
        textbox(s, x + 0.26, 4.27, 2.9, 0.88, body, size=10.5, color=MUTED)
        textbox(s, x + 0.26, 5.28, 2.9, 0.27, result, size=11, color=c, bold=True)
    textbox(s, 0.76, 6.17, 11.6, 0.5, "The compiler evaluates primitive facts and project policy—not a universal “architecture score.”", size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "All three frozen semantic candidates carry accepted six-of-six behavior results supplied at the transaction boundary. The compiler evaluates their validated semantic graphs against locked policy. The domain-owned version publishes. The centralized version is rolled back for dispatch growth and transport authority, state, and dependency violations. A superficial helper split is also rolled back because aggregate responsibility remains in transport. This is not yet a complete source-to-agent product loop.")


def slide_09(prs):
    s = base_slide(prs, 9, "ATOMIC CONTROL")
    title(s, "The candidate revision is the unit of trust", "Every proof binds to the same immutable revision.")
    nodes = [
        (0.7, 2.65, "BASE", "immutable revision", BLUE),
        (3.15, 2.65, "CANDIDATE", "whole workspace", CYAN),
        (5.6, 2.65, "VALIDATE", "static • impact • behavior • policy", CYAN),
    ]
    for x, y, h, b, c in nodes:
        rect(s, x, y, 2.05, 1.45, PANEL, c)
        textbox(s, x + 0.15, y + 0.28, 1.75, 0.3, h, size=14, color=c, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, x + 0.15, y + 0.82, 1.75, 0.28, b, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    line(s, 2.78, 3.38, 3.05, 3.38, BLUE, 2, True)
    line(s, 5.23, 3.38, 5.5, 3.38, CYAN, 2, True)
    # fork
    line(s, 7.67, 3.05, 8.25, 2.25, GREEN, 2, True)
    line(s, 7.67, 3.75, 8.25, 4.72, RED, 2, True)
    rect(s, 8.42, 1.65, 3.55, 1.55, PANEL, GREEN)
    textbox(s, 8.68, 1.95, 3.0, 0.35, "PUBLISH CHILD", size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 8.68, 2.48, 3.0, 0.28, "completion evidence attached", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    rect(s, 8.42, 4.18, 3.55, 1.55, PANEL, RED)
    textbox(s, 8.68, 4.48, 3.0, 0.35, "ROLL BACK", size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 8.68, 5.01, 3.0, 0.28, "base remains unchanged", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    textbox(s, 0.78, 5.2, 6.2, 0.8, "No partial publication.\nNo stale completion report.\nNo “tests passed” detached from the final code.", size=15, color=WHITE, bold=True)
    add_notes(s, "AIL’s revision protocol matters commercially because evidence cannot drift away from the code it describes. A whole candidate is canonicalized and checked before publication. A successful change publishes one child revision with bound completion evidence. A denied or incomplete change publishes nothing and leaves the base intact.")


def slide_10(prs):
    s = base_slide(prs, 10, "MODEL INDEPENDENCE")
    title(s, "Models improve. The control layer compounds.", "AIL’s semantic interface is designed to serve any capable coding agent.")
    models = [(0.85, "FRONTIER\nMODEL"), (0.85, "OPEN\nMODEL"), (0.85, "SPECIALIZED\nMODEL")]
    ys = [2.15, 3.52, 4.89]
    for (_, label), y in zip(models, ys):
        rect(s, 0.85, y, 2.15, 0.9, PANEL, BLUE)
        textbox(s, 0.98, y + 0.2, 1.88, 0.43, label, size=12, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        line(s, 3.05, y + 0.45, 4.4, 3.97, BLUE, 1.5, True)
    rect(s, 4.52, 2.42, 3.4, 3.12, PANEL_2, CYAN)
    textbox(s, 4.85, 2.82, 2.75, 0.45, "AIL SEMANTIC API", size=19, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 4.9, 3.63, 2.65, 1.2, "revision-bound facts\ndeterministic queries\nvalidated operations", size=13, align=PP_ALIGN.CENTER)
    outputs = [(9.42, 2.22, "VALIDATED OPERATION"), (9.42, 3.57, "REVIEW EVIDENCE"), (9.42, 4.92, "POLICY RESULT")]
    for x, y, label in outputs:
        line(s, 7.98, 3.97, 9.25, y + 0.42, CYAN, 1.5, True)
        rect(s, x, y, 2.45, 0.84, PANEL, GREEN)
        textbox(s, x + 0.12, y + 0.24, 2.2, 0.26, label, size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 0.84, 6.32, 11.2, 0.34, "Not a prompt format. Not tied to one model vendor. Canonical source stays portable and auditable.", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(s, "AIL is intentionally not a prompt language or a model-specific encoding. Better models should improve the agent layer, while AIL’s accumulated semantic contracts, compiler implementation, fixtures, and project policy remain reusable. This creates a durable control surface across model cycles and lets teams choose models based on cost, privacy, or capability.")


def slide_11(prs):
    s = base_slide(prs, 11, "ECONOMIC HYPOTHESIS")
    title(s, "The leverage is in the total cost of a correct change", "Compiler mechanisms execute as specified. Comparative advantage remains to be measured.")
    terms = [
        ("GENERATION", BLUE, "already falling"),
        ("CONTEXT", CYAN, "query, don’t rediscover"),
        ("CONSEQUENCE", CYAN, "exact impact + authority"),
        ("REPAIR", GREEN, "structured causes"),
        ("REGRESSION", GREEN, "policy before publish"),
        ("REVIEW", AMBER, "bound evidence"),
    ]
    for i, (name, c, sub) in enumerate(terms):
        x = 0.65 + i * 2.08
        height = 1.9
        y = 5.25 - height
        rect(s, x, y, 1.65, height, c, c, radius=False)
        textbox(s, x + 0.08, y + 0.17, 1.49, 0.32, name, size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, x + 0.13, y + height - 0.66, 1.39, 0.5, sub, size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    line(s, 0.65, 5.42, 12.1, 5.42, MUTED, 1)
    textbox(s, 0.69, 5.72, 11.4, 0.42, "HYPOTHESIS: lower tokens + elapsed time + retries + regressions + reviewer effort at equal correctness", size=13, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 0.69, 6.3, 11.4, 0.3, "No speedup or ROI number is claimed today.", size=11, color=RED, align=PP_ALIGN.CENTER)
    textbox(s, 0.72, 2.98, 4.0, 0.24, "COMPONENTS SHOWN EQUALLY — NOT TO SCALE", size=8.5, color=AMBER, bold=True)
    add_notes(s, "The investment case is not fewer source tokens. It is lower total work to reach a correct, reviewable change. We will measure model tokens, elapsed work, repair cycles, regressions, and reviewer effort under equal correctness gates. We do not yet have a defensible speedup or ROI number, and the deck should say that plainly.")


def slide_12(prs):
    s = base_slide(prs, 12, "BEACHHEAD + DEFENSIBILITY")
    title(s, "Start where agent autonomy is already the operating model", "Buyer and adoption path are hypotheses to validate with design partners.")
    card(s, 0.7, 2.0, 3.7, 3.8, "BEACHHEAD", "Greenfield backend services and workers\n\nBuyer hypothesis:\nEngineering leaders running agent-primary delivery\n\nEntry point:\nCompiler-guided validation in bounded service pilots", accent=BLUE, body_size=12)
    card(s, 4.82, 2.0, 3.7, 3.8, "WHY A LANGUAGE", "Existing-language tools can improve retrieval and analysis.\n\nAIL is justified only where canonical semantics, explicit authority, deterministic behavior, and complete supported impact cannot be added reliably through compatibility-bound tooling.", accent=CYAN, body_size=11.5)
    card(s, 8.94, 2.0, 3.7, 3.8, "HOW DEFENSIBILITY CAN COMPOUND", "Built: compiler • contracts • semantic graph • adversarial fixtures\n\nTo build: held-out workloads • design-partner policy • integrations • production runtime • adoption", accent=GREEN, body_size=11.5)
    textbox(s, 0.8, 6.19, 11.7, 0.4, "No customer adoption, proprietary workload access, or switching-cost claim is made today.", size=11, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "The first buyer hypothesis is an engineering organization using agents as primary implementers for greenfield backend services and workers. The funded entry point is a bounded validation pilot, not production deployment. AIL belongs in a language only where foundational guarantees—canonical semantics, explicit authority, determinism, and complete supported impact—cannot be reliably retrofitted through compatibility-bound tooling. Defensibility is a strategy: the compiler, contracts, graph, and fixtures exist; held-out workloads, design-partner policy, integrations, production runtime, and adoption remain to build.")


def slide_13(prs):
    s = base_slide(prs, 13, "ROADMAP")
    title(s, "Mechanism first. Economics next. Production only when justified.")
    phases = [
        (0.7, "DONE", "Working compiler", "semantic inspection\nchange transactions\nM28 calls + modules", GREEN),
        (3.8, "TEST", "Comparative validation", "fresh held-out tasks\nrequired AIL capability\nstrong baseline trials", CYAN),
        (6.9, "DECIDE", "Use the results", "correctness first\ntotal change cost\ngeneralization", BLUE),
        (10.0, "IF NEEDED", "Broader capability", "language, runtime,\nlowering, or ecosystem\nwork evidence requires", AMBER),
    ]
    for i, (x, phase, h, b, c) in enumerate(phases):
        pill(s, x, 2.0, 1.1, phase, c, NAVY)
        rect(s, x, 2.58, 2.65, 2.85, PANEL, c)
        textbox(s, x + 0.2, 2.87, 2.25, 0.55, h, size=17, color=c, bold=True)
        textbox(s, x + 0.2, 3.7, 2.25, 1.15, b, size=12, color=MUTED)
        if i < 3:
            line(s, x + 2.68, 4.0, x + 3.0, 4.0, c, 2, True)
    textbox(s, 0.75, 5.94, 11.6, 0.6, "M28 is complete. No successor is active; this comparison is one possible next investment.", size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "M28 shipped modules, import aliases, qualified references, and ordinary calls. The repository has no active successor milestone. One possible next investment is to freeze representative tasks, build only the AIL capabilities needed for a fair comparator, and compare against mainstream languages with their normal strong tools. Native lowering, LLVM integration, broader syntax, and self-hosting are not automatic next steps.")


def slide_14(prs):
    s = base_slide(prs, 14, "CURRENT LIMITS")
    title(s, "What works—and what remains to build or measure")
    rect(s, 0.7, 1.95, 5.72, 4.55, PANEL, GREEN)
    pill(s, 0.98, 2.23, 1.45, "Proven today", GREEN, NAVY)
    proven = [
        "Executable Rust compiler, not a UI mock",
        "Deterministic parse, check, execute, and revisions",
        "Exact schema-impact report for the M19 contract",
        "Atomic whole-candidate validation",
        "Architecture policy blocks behavior-passing regressions",
        "37-case public AIL behavior corpus passes",
    ]
    for i, item in enumerate(proven):
        dot(s, 1.02, 3.02 + i * 0.5, 0.11, GREEN)
        textbox(s, 1.28, 2.93 + i * 0.5, 4.7, 0.35, item, size=11.5)
    rect(s, 6.88, 1.95, 5.72, 4.55, PANEL, AMBER)
    pill(s, 7.16, 2.23, 1.72, "To prove next", AMBER, NAVY)
    unproven = [
        "Lower total agent-change cost vs strong baselines",
        "Repeatability across models and representative tasks",
        "Lower regression and repair rates",
        "Lower human reviewer effort",
        "Production runtime and native lowering",
        "Broad language and ecosystem coverage",
    ]
    for i, item in enumerate(unproven):
        dot(s, 7.2, 3.02 + i * 0.5, 0.11, AMBER)
        textbox(s, 7.46, 2.93 + i * 0.5, 4.7, 0.35, item, size=11.5)
    textbox(s, 7.45, 6.03, 4.7, 0.24, "AIL-VS-BASELINE AGENT TRIALS: 0", size=10, color=RED, bold=True)
    add_notes(s, "The compiler provides exact semantic impact for the M19 contract and atomic enforcement of the M24 architecture rules. No comparison has yet shown lower cost than Rust, Go, Python, or TypeScript with their normal tools. AIL also lacks native lowering, production I/O, general concurrency, and broad ecosystem support.")


def slide_15(prs):
    s = base_slide(prs, 15, "NEXT EXPERIMENT")
    title(s, "Run the comparison that can falsify the thesis", "Measure the result instead of assuming a speedup.")
    plans = [
        (0.7, "1", "LOCK", "Fresh held-out design-partner tasks\nEquivalent Rust / Go / Python / TS starts\nFixed max attempts, tools, oracles, review rubric", BLUE),
        (4.47, "2", "RUN", "Mechanism fixtures calibrate the pipeline\nBaseline targets lock before AIL trials\nCompletion, tokens, time, repairs, regressions", CYAN),
        (8.24, "3", "DECIDE", "Independent randomized review\nUncertainty + failure distributions\nChoose the next build from results\nProduction decision", GREEN),
    ]
    for x, n, h, b, c in plans:
        rect(s, x, 2.08, 3.42, 3.5, PANEL, c)
        textbox(s, x + 0.25, 2.34, 0.55, 0.55, n, size=29, color=c, bold=True)
        textbox(s, x + 0.88, 2.46, 2.2, 0.32, h, size=15, color=c, bold=True)
        textbox(s, x + 0.25, 3.22, 2.9, 1.75, b, size=11.5, color=MUTED)
    textbox(s, 0.78, 5.95, 11.55, 0.58, "Fund: compiler/runtime engineering  •  agent-evaluation infrastructure  •  design-partner workload access", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "First, use UC-003 and CancelJob to test the runner, then freeze fresh held-out tasks before adapting AIL to them. Build equivalent mainstream-language workspaces and fix the maximum attempts, outcomes, exclusions, and reviewer rubric before running. Second, establish baseline distributions and lock AIL success targets before AIL trials. Third, run randomized independent review and choose the next build from the results. Financing, runway, and team requirements are not specified here.")


def slide_16(prs):
    s = base_slide(prs, 16, "APPENDIX: LIVE DEMO")
    title(s, "Retained M27 pilot: an operator changes course", "The seeded bad candidate was supplied; the recorded run demonstrates one repair, not comparative advantage.")
    flow = [
        (0.6, "INPUT", "seeded centralized\ncandidate", BLUE),
        (2.67, "REJECTION", "6/6 behavior\npublication denied", RED),
        (4.74, "DRILL-DOWN", "exact policy\ncontributors", CYAN),
        (6.81, "REPAIR", "move authority\nto domain", CYAN),
        (8.88, "VALIDATE", "6/6 behavior\n0 findings", GREEN),
        (10.95, "EVIDENCE", "publish child\nreview bundle", GREEN),
    ]
    for i, (x, h, b, c) in enumerate(flow):
        rect(s, x, 2.25, 1.75, 1.42, PANEL, c)
        textbox(s, x + 0.12, 2.47, 1.5, 0.27, h, size=10, color=c, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, x + 0.12, 2.9, 1.5, 0.46, b, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 5:
            line(s, x + 1.78, 2.96, x + 2.0, 2.96, c, 1.8, True)
    left = ["Timer + cumulative model input tokens", "Source reads + semantic queries", "Validation attempts + repair cycles"]
    right = ["Public + hidden regressions", "Final revision + compiler evidence", "Blind reviewer time, questions, defects"]
    card(s, 0.78, 4.35, 5.72, 1.72, "AGENT TELEMETRY", "\n".join(f"• {x}" for x in left), accent=CYAN, body_size=11)
    card(s, 6.82, 4.35, 5.72, 1.72, "OUTCOME + REVIEW", "\n".join(f"• {x}" for x in right), accent=GREEN, body_size=11)
    textbox(s, 0.8, 6.33, 11.5, 0.32, "Do not compare a strong AIL workflow with a deliberately crippled baseline.", size=11, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "M27 retained one complete run. The operator received the seeded centralized candidate and its compact rejection; the run does not show the operator independently creating that bad candidate. The operator used structured contributors, moved authority to the domain handler, reran validation, and published the valid child. A comparison must give baseline languages their normal compilers, language servers, refactors, search, formatters, and tests.")


def slide_17(prs):
    s = base_slide(prs, 17, "APPENDIX: MEASUREMENT")
    title(s, "Baseline vs AIL: one task contract, two strong toolchains", "Completion and terminal failures are primary; successful-run efficiency is conditional.")
    # two lanes
    for x, label, c, tools in [
        (0.7, "MAINSTREAM BASELINE", BLUE, "Rust / Go / Python / TypeScript\ncompiler + LSP + search + formatter + tests"),
        (6.87, "AIL TREATMENT", CYAN, "same model + agent policy\nAIL compiler semantic queries + validated operations"),
    ]:
        rect(s, x, 1.95, 5.72, 1.55, PANEL, c)
        textbox(s, x + 0.26, 2.23, 5.2, 0.3, label, size=13, color=c, bold=True)
        textbox(s, x + 0.26, 2.72, 5.2, 0.45, tools, size=11, color=MUTED)
    line(s, 6.64, 1.98, 6.64, 6.32, PANEL_2, 2)
    dimensions = [
        ("WORK", "elapsed agent time • tool calls • source reads"),
        ("TOKENS", "provider-counted input by category • repeated context"),
        ("REPAIR", "validation attempts • repair cycles • terminal failures"),
        ("QUALITY", "public/private checks • seeded misses • new authority"),
        ("REVIEW", "independent reviewer time • questions • defects • confidence"),
    ]
    for i, (h, b) in enumerate(dimensions):
        y = 3.84 + i * 0.52
        textbox(s, 0.92, y, 1.35, 0.28, h, size=10, color=CYAN, bold=True)
        textbox(s, 2.25, y, 9.8, 0.28, b, size=10.5, color=WHITE)
    textbox(s, 0.76, 6.52, 11.7, 0.28, "Fixed max attempts • pre-registered summaries + uncertainty • targets locked before AIL results", size=10.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "The comparison uses the same task contract and correctness oracle. Baselines receive strong normal tooling. Use the same model and agent policy where possible, and record version and permission differences. Fix the maximum attempts per cell before starting. Completion and terminal-failure distributions are primary outcomes; token, time, and repair comparisons among successful runs are conditional secondary measurements. Randomize reviewer assignments. Reviewers can be blinded to the hypothesis and provenance where feasible, but not to a language that is visible in source.")


def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "AIL — A programming language for software built by agents"
    prs.core_properties.subject = "Investor presentation and current engineering status"
    prs.core_properties.author = "AIL"
    prs.core_properties.keywords = "AIL, agents, compiler, investor, programming language"

    for builder in (
        slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
        slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
        slide_13, slide_14, slide_15, slide_16, slide_17,
    ):
        builder(prs)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Wrote {output} ({len(prs.slides)} slides)")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    build(args.output.resolve())


if __name__ == "__main__":
    main()
